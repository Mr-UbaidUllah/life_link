#!/usr/bin/env python3
"""Generate the Life Link FYP presentation (.pptx)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# ---- Palette --------------------------------------------------------------
RED      = RGBColor(0xC6, 0x28, 0x28)   # primary crimson
DARK_RED = RGBColor(0x8E, 0x00, 0x00)
INK      = RGBColor(0x21, 0x25, 0x2B)   # near-black text
SLATE    = RGBColor(0x55, 0x5C, 0x66)   # secondary text
LIGHT    = RGBColor(0xF6, 0xF7, 0xF9)   # light panel
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT   = RGBColor(0xE5, 0x39, 0x35)   # accent red

IMG_DIR = "diagrams/images"
EMU_IN = 914400

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def set_run(run, text, size, color, bold=False, italic=False, font="Calibri"):
    run.text = text
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def header(slide, kicker, title):
    """Standard content-slide header: accent bar + kicker + title."""
    rect(slide, 0, 0, Inches(0.28), SH, RED)            # left spine
    rect(slide, Inches(0.6), Inches(0.55), Inches(0.55), Inches(0.09), ACCENT)
    tb, tf = textbox(slide, Inches(0.6), Inches(0.68), Inches(12), Inches(0.4))
    set_run(tf.paragraphs[0].add_run(), kicker.upper(), 12, RED, bold=True)
    tb, tf = textbox(slide, Inches(0.58), Inches(1.0), Inches(12.2), Inches(0.9))
    set_run(tf.paragraphs[0].add_run(), title, 30, INK, bold=True)
    return Inches(2.0)  # content top


def bullets(slide, items, top=Inches(2.05), left=Inches(0.75),
            width=Inches(11.8), size=17, gap=10):
    tb, tf = textbox(slide, left, top, width, SH - top - Inches(0.5))
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            text, lvl = it
        else:
            text, lvl = it, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0)
        p.level = lvl
        bullet = "•   " if lvl == 0 else "–   "
        r = p.add_run()
        set_run(r, bullet + text, size - (lvl * 1), INK if lvl == 0 else SLATE,
                bold=(lvl == 0 and text.endswith(":")))
    return tb


def fit_image(path, box_x, box_y, box_w, box_h):
    """Return (x,y,w,h) to fit image inside box, centered."""
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    bar = box_w / box_h
    if ar > bar:
        w = box_w; h = int(box_w / ar)
    else:
        h = box_h; w = int(box_h * ar)
    x = box_x + (box_w - w) // 2
    y = box_y + (box_h - h) // 2
    return x, y, w, h


def diagram_slide(kicker, title, img, caption, note):
    s = add_slide()
    rect(s, 0, 0, SW, SH, WHITE)
    header(s, kicker, title)
    # image panel
    box_x, box_y = Inches(0.7), Inches(1.95)
    box_w, box_h = Inches(11.9), Inches(4.55)
    rect(s, box_x, box_y, box_w, box_h, LIGHT)
    x, y, w, h = fit_image(os.path.join(IMG_DIR, img),
                           box_x + Inches(0.1), box_y + Inches(0.1),
                           box_w - Inches(0.2), box_h - Inches(0.2))
    s.shapes.add_picture(os.path.join(IMG_DIR, img), x, y, width=w, height=h)
    tb, tf = textbox(s, Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.6))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_run(tf.paragraphs[0].add_run(), caption, 13, SLATE, italic=True)
    notes(s, note)
    return s


def chip(slide, x, y, w, label, value, vcolor=RED):
    h = Inches(1.25)
    card = rect(slide, x, y, w, h, LIGHT)
    tb, tf = textbox(slide, x, y + Inches(0.16), w, Inches(0.6),
                     anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_run(tf.paragraphs[0].add_run(), value, 26, vcolor, bold=True)
    tb, tf = textbox(slide, x, y + Inches(0.74), w, Inches(0.45),
                     anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_run(tf.paragraphs[0].add_run(), label, 12, SLATE, bold=True)


# ===========================================================================
# SLIDE 1 — TITLE
# ===========================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, 0, Inches(0.45), SH, RED)
rect(s, 0, SH - Inches(0.18), SW, Inches(0.18), RED)
# university logo (top-right)
LOGO = "assets_pres/uni_logo.png"
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(10.55), Inches(0.55), width=Inches(2.0), height=Inches(2.0))
# brand
tb, tf = textbox(s, Inches(1.0), Inches(1.25), Inches(9), Inches(0.5))
set_run(tf.paragraphs[0].add_run(), "FINAL YEAR PROJECT II", 15, ACCENT, bold=True)
tb, tf = textbox(s, Inches(0.95), Inches(1.78), Inches(9.4), Inches(1.5))
set_run(tf.paragraphs[0].add_run(), "Life Link", 66, WHITE, bold=True)
tb, tf = textbox(s, Inches(1.0), Inches(3.05), Inches(11), Inches(0.8))
set_run(tf.paragraphs[0].add_run(),
        "A Blood Donation & Emergency Response Mobile Platform",
        23, RGBColor(0xE8, 0xE8, 0xEA))
tb, tf = textbox(s, Inches(1.0), Inches(3.82), Inches(11), Inches(0.6))
set_run(tf.paragraphs[0].add_run(),
        "Connecting blood donors, patients and emergency services in real time",
        14, SLATE, italic=True)
# group members band
tb, tf = textbox(s, Inches(1.0), Inches(4.75), Inches(11.3), Inches(1.0))
p = tf.paragraphs[0]
set_run(p.add_run(), "Presented by:   ", 14, ACCENT, bold=True)
set_run(p.add_run(), "Ubaid Ullah (179)      M. Saad Awan (215)      Umar Rashid (427)",
        14, WHITE, bold=True)
# supervisor + degree
tb, tf = textbox(s, Inches(1.0), Inches(5.55), Inches(11.3), Inches(1.4))
p = tf.paragraphs[0]
set_run(p.add_run(), "Supervised by:  ", 13, SLATE, bold=True)
set_run(p.add_run(), "Engr. Waseem Ullah Khan  (Lecturer, ICS/IT)", 13, WHITE)
p2 = tf.add_paragraph(); p2.space_before = Pt(5)
set_run(p2.add_run(), "Bachelor of Science in Computer Science   •   Session 2022–2026",
        13, RGBColor(0xC9, 0xCC, 0xD1))
p3 = tf.add_paragraph(); p3.space_before = Pt(5)
set_run(p3.add_run(),
        "Institute of Computer Sciences & IT (ICS/IT)  •  The University of Agriculture, Peshawar",
        12, SLATE)
notes(s, "Good morning/afternoon respected panel members and supervisor. We are Ubaid Ullah, "
         "Muhammad Saad Awan and Umar Rashid, and today we present our Final Year Project II, "
         "'Life Link' — a mobile platform built to solve the problem of finding blood donors quickly "
         "during emergencies. The project is supervised by Engr. Waseem Ullah Khan and submitted to "
         "the Institute of Computer Sciences and IT at The University of Agriculture, Peshawar. "
         "Over the next ~15 minutes we'll cover the problem, our proposed solution, the features we "
         "built, the technology and architecture, and the design diagrams.")

# ===========================================================================
# SLIDE 2 — AGENDA
# ===========================================================================
s = add_slide(); rect(s, 0, 0, SW, SH, WHITE)
header(s, "Overview", "Presentation Agenda")
left_items = ["Problem statement & motivation",
              "Existing vs. proposed system",
              "Aim, objectives & scope",
              "Key features of Life Link",
              "Technology stack"]
right_items = ["System architecture & design",
               "UML / DFD / ERD diagrams",
               "Database & security model",
               "Testing & results",
               "Challenges, future work & conclusion"]
tb = bullets(s, [f"{i+1}.  {t}" for i, t in enumerate(left_items)],
             top=Inches(2.3), left=Inches(0.9), width=Inches(5.8), size=18, gap=16)
tb = bullets(s, [f"{i+6}.  {t}" for i, t in enumerate(right_items)],
             top=Inches(2.3), left=Inches(7.0), width=Inches(5.8), size=18, gap=16)
notes(s, "Here is the roadmap for my presentation. I'll start with the problem and motivation, "
         "compare the current manual process with my proposed app, state my objectives, then walk "
         "through the features, the technology, and the engineering diagrams. I'll finish with "
         "testing, challenges and future work. Please feel free to ask questions at the end.")

# ===========================================================================
# SLIDE 3 — PROBLEM STATEMENT
# ===========================================================================
s = add_slide(); rect(s, 0, 0, SW, SH, WHITE)
header(s, "Motivation", "Problem Statement")
bullets(s, [
    "Every year thousands of patients need urgent blood transfusions, yet finding a matching donor in time is difficult.",
    "Blood requests today spread through phone calls, WhatsApp groups and word of mouth — slow, unreliable and limited in reach.",
    ("There is no single platform that instantly connects a patient to nearby, blood-group-matched, willing donors.", 1),
    ("Donors who are willing to help have no easy way to discover requests near them.", 1),
    "Emergency contacts — ambulances, blood banks, volunteers — are scattered and hard to reach in a crisis.",
    "Result: critical time is lost, and avoidable deaths occur due to delayed blood availability.",
], top=Inches(2.1), size=17, gap=12)
notes(s, "Let me start with the core problem. In an emergency, finding the right blood at the right "
         "time is genuinely hard. Today people rely on phone calls and WhatsApp groups, which are "
         "slow and reach only a small circle. There's no central system that matches a patient with "
         "nearby donors of the correct blood group, and willing donors have no easy way to find "
         "requests around them. On top of that, emergency contacts like ambulances and blood banks "
         "are scattered. The cost of this delay can literally be a life. This is the gap Life Link "
         "is designed to close.")

# ===========================================================================
# SLIDE 4 — EXISTING SYSTEM (diagram)
# ===========================================================================
diagram_slide("Current Process", "Existing System — The Manual Way",
    "05_existing_system.png",
    "Figure 1: How blood is sourced today — manual calls, scattered contacts, no matching.",
    "This diagram shows how the process works today. A patient or relative manually calls hospitals, "
    "relatives and friends, posts in WhatsApp groups, and hopes someone responds. There is no "
    "filtering by blood group or location, no tracking, and lots of wasted time. The key weaknesses "
    "are: it's slow, it's manual, it has limited reach, and there's no guarantee of a matching donor.")

# ===========================================================================
# SLIDE 5 — PROPOSED SYSTEM (diagram)
# ===========================================================================
diagram_slide("Our Solution", "Proposed System — Life Link",
    "06_proposed_system.png",
    "Figure 2: Life Link automates matching, notifies donors instantly, and centralizes services.",
    "This is my proposed solution. With Life Link, a patient posts a blood request once. The system "
    "automatically finds donors with the matching blood group in the same city and instantly pushes a "
    "notification to them. Donors can respond and chat in real time. All emergency services — "
    "ambulances, organizations and volunteers — live in one app. So the process becomes automated, "
    "targeted, fast, and far wider in reach than manual calling.")

# ===========================================================================
# SLIDE 6 — AIM & OBJECTIVES
# ===========================================================================
s = add_slide(); rect(s, 0, 0, SW, SH, WHITE)
header(s, "Goals", "Aim & Objectives")
tb, tf = textbox(s, Inches(0.75), Inches(2.0), Inches(11.8), Inches(0.9))
p = tf.paragraphs[0]
set_run(p.add_run(), "Aim:  ", 18, RED, bold=True)
set_run(p.add_run(), "To build a mobile platform that connects blood donors and patients quickly "
        "and reliably, reducing the time it takes to find a matching donor in an emergency.",
        18, INK)
bullets(s, [
    "Objectives:",
    ("Let users register, build a donor profile and toggle their availability.", 1),
    ("Allow patients to post blood requests with blood group, hospital and location.", 1),
    ("Automatically match requests to suitable donors and notify them instantly.", 1),
    ("Provide search to find donors by blood group and location.", 1),
    ("Offer a directory of ambulances, blood banks and volunteers.", 1),
    ("Enable real-time, in-app chat between donors and patients.", 1),
    ("Secure all data with authentication and access rules.", 1),
], top=Inches(3.0), size=16, gap=9)
notes(s, "My overall aim was to reduce the time it takes to find a matching donor. To achieve that I "
         "set seven concrete objectives: user and donor registration, posting blood requests, automatic "
         "matching with instant notifications, donor search, an emergency services directory, real-time "
         "chat, and strong data security. Each objective maps directly to a feature I'll show you next, "
         "so the project is measurable against these goals.")

# ===========================================================================
# SLIDE 7 — SCOPE
# ===========================================================================
s = add_slide(); rect(s, 0, 0, SW, SH, WHITE)
header(s, "Boundaries", "Project Scope")
# Two columns: In scope / Out of scope
rect(s, Inches(0.75), Inches(2.1), Inches(5.7), Inches(0.6), RGBColor(0x1B,0x5E,0x20))
tb, tf = textbox(s, Inches(0.75), Inches(2.1), Inches(5.7), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
set_run(tf.paragraphs[0].add_run(), "IN SCOPE", 16, WHITE, bold=True)
bullets(s, [
    "Mobile app (Android & iOS) built with Flutter",
    "Donor & patient accounts via Firebase Auth",
    "Blood request posting + automatic donor matching",
    "Push notifications for matched requests",
    "Real-time chat between users",
    "Emergency services directory",
], top=Inches(2.85), left=Inches(0.85), width=Inches(5.5), size=14, gap=9)

rect(s, Inches(6.9), Inches(2.1), Inches(5.7), Inches(0.6), SLATE)
tb, tf = textbox(s, Inches(6.9), Inches(2.1), Inches(5.7), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
set_run(tf.paragraphs[0].add_run(), "OUT OF SCOPE", 16, WHITE, bold=True)
bullets(s, [
    "Online payments / monetary transactions",
    "Medical verification of blood compatibility by labs",
    "GPS turn-by-turn navigation to hospitals",
    "Integration with government blood-bank databases",
    "Web/desktop version (mobile-first project)",
], top=Inches(2.85), left=Inches(7.0), width=Inches(5.5), size=14, gap=9)
notes(s, "It's important in an FYP to be clear about boundaries. In scope: a cross-platform mobile app "
         "with donor and patient accounts, request posting, automatic matching, notifications, chat and a "
         "services directory. Out of scope: payments, lab-level medical verification, live GPS navigation, "
         "government database integration, and a web version. Defining this kept the project achievable in "
         "the available time and is a natural lead-in to my future work slide.")

# ===========================================================================
# SLIDE 8 — KEY FEATURES OVERVIEW
# ===========================================================================
s = add_slide(); rect(s, 0, 0, SW, SH, WHITE)
header(s, "Product", "Key Features at a Glance")
feats = [
    ("Donor Profiles", "Register, set blood group & toggle availability"),
    ("Blood Requests", "Post urgent needs with hospital & location"),
    ("Smart Matching", "Auto-match by blood group + city"),
    ("Find Donors", "Search & filter willing donors nearby"),
    ("Real-time Chat", "Message donors with read receipts"),
    ("Push Alerts", "Instant notifications to matched donors"),
    ("Ambulances", "Directory of emergency ambulance services"),
    ("Organizations", "Blood banks & partner NGOs"),
    ("Volunteers", "Community volunteer directory"),
]
cols, cw, ch = 3, Inches(3.85), Inches(1.35)
gx, gy = Inches(0.75), Inches(2.15)
for i, (t, d) in enumerate(feats):
    r, c = divmod(i, cols)
    x = gx + c * (cw + Inches(0.2))
    y = gy + r * (ch + Inches(0.18))
    rect(s, x, y, cw, ch, LIGHT)
    rect(s, x, y, Inches(0.1), ch, RED)
    tb, tf = textbox(s, x + Inches(0.25), y + Inches(0.14), cw - Inches(0.35), Inches(0.45))
    set_run(tf.paragraphs[0].add_run(), t, 15, RED, bold=True)
    tb, tf = textbox(s, x + Inches(0.25), y + Inches(0.6), cw - Inches(0.35), Inches(0.65))
    set_run(tf.paragraphs[0].add_run(), d, 12, SLATE)
notes(s, "Here's the whole product on one slide — nine core features grouped into three areas. "
         "First, donation: donor profiles, blood requests and smart matching. Second, discovery and "
         "communication: find donors, real-time chat and push alerts. Third, the emergency directory: "
         "ambulances, organizations and volunteers. I'll now go a little deeper into the most important "
         "ones over the next few slides.")

# ===========================================================================
# SLIDE 9 — FEATURE: REQUESTS & MATCHING
# ===========================================================================
s = add_slide(); rect(s, 0, 0, SW, SH, WHITE)
header(s, "Core Feature", "Blood Requests & Smart Matching")
bullets(s, [
    "A patient posts a request: blood group, number of bags, hospital, reason, contact and location.",
    "On creation, a Firebase Cloud Function fires automatically — no manual work.",
    ("It queries all users, keeps only those whose blood group AND city match the request.", 1),
    ("Matching is case-insensitive and excludes the request creator.", 1),
    ("Matched donors receive an instant push notification via Firebase Cloud Messaging.", 1),
    "Requests carry a status (open/closed) and an expiry date, so stale requests disappear.",
    "This turns a slow manual phone-tree into an instant, targeted broadcast.",
], top=Inches(2.05), size=16, gap=11)
notes(s, "This is the heart of the app. When a patient submits a blood request, a Cloud Function on "
         "the backend triggers automatically. It looks through the user base and selects only the donors "
         "whose blood group and city match — comparison is case-insensitive so 'a+' and 'A+' both work — "
         "and it never notifies the person who made the request. Those matched donors get an instant push "
         "notification. Requests also have an open/closed status and an expiry date so the feed stays "
         "fresh. The big win: a process that used to take dozens of phone calls now happens in seconds.")

# ===========================================================================
# SLIDE 10 — FEATURE: DONORS, CHAT, DIRECTORY
# ===========================================================================
s = add_slide(); rect(s, 0, 0, SW, SH, WHITE)
header(s, "Core Features", "Donors, Chat & Emergency Directory")
bullets(s, [
    "Be a Donor: toggle availability; turning it on notifies the community that a new donor is available.",
    "Find Donors: search and filter the donor list by blood group and location.",
    "Real-time Chat: one-to-one messaging built on Firestore streams.",
    ("Unread counts, read receipts and last-message previews — like a modern chat app.", 1),
    "Emergency Directory — three curated lists for a crisis:",
    ("Ambulances — name, hospital, address and one-tap call.", 1),
    ("Organizations — blood banks and partner NGOs.", 1),
    ("Volunteers — community helpers with their work description.", 1),
], top=Inches(2.05), size=16, gap=10)
notes(s, "Beyond requests, donors can flip an availability switch, and when they do, the community is "
         "notified that a new donor is available. Anyone can search and filter the donor list by blood "
         "group and city. For coordination, I built real-time chat on Firestore streams, with unread "
         "counts, read receipts and message previews — it feels like a normal messaging app. Finally, the "
         "emergency directory gives three ready-to-use lists — ambulances, organizations and volunteers — "
         "each with one-tap calling, so help is never more than a tap away.")

# ===========================================================================
# SLIDE 11 — TECH STACK
# ===========================================================================
s = add_slide(); rect(s, 0, 0, SW, SH, WHITE)
header(s, "Engineering", "Technology Stack")
groups = [
    ("FRONTEND", ["Flutter (Dart) — cross-platform", "Provider — state management",
                  "flutter_screenutil — responsive UI", "google_fonts, flutter_svg, shimmer"]),
    ("BACKEND", ["Firebase Authentication", "Cloud Firestore (NoSQL DB)",
                 "Firebase Storage (images)", "Cloud Functions (Node.js)"]),
    ("SERVICES", ["Firebase Cloud Messaging (FCM)", "flutter_local_notifications",
                  "image_picker", "url_launcher (calls/links)"]),
]
cw = Inches(3.95)
for i, (title, items) in enumerate(groups):
    x = Inches(0.75) + i * (cw + Inches(0.2))
    rect(s, x, Inches(2.15), cw, Inches(0.62), RED)
    tb, tf = textbox(s, x, Inches(2.15), cw, Inches(0.62), anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_run(tf.paragraphs[0].add_run(), title, 16, WHITE, bold=True)
    rect(s, x, Inches(2.77), cw, Inches(3.4), LIGHT)
    bullets(s, items, top=Inches(3.0), left=x + Inches(0.2),
            width=cw - Inches(0.35), size=14, gap=12)
tb, tf = textbox(s, Inches(0.75), Inches(6.45), Inches(11.8), Inches(0.6))
set_run(tf.paragraphs[0].add_run(),
        "Why Flutter + Firebase?  One codebase for Android & iOS, real-time database out of the box, "
        "and a serverless backend that scales without managing servers.", 13, SLATE, italic=True)
notes(s, "On the engineering side, the front end is Flutter with Dart, using the Provider pattern for "
         "state management and ScreenUtil for responsive layouts. The backend is entirely Firebase: "
         "Authentication for login, Cloud Firestore as the real-time NoSQL database, Storage for profile "
         "images, and Cloud Functions in Node.js for server logic. Notifications use Firebase Cloud "
         "Messaging. I chose Flutter plus Firebase because one codebase ships to both Android and iOS, "
         "the database is real-time by default, and the backend is serverless — so as a single student I "
         "didn't have to manage any servers.")

# ===========================================================================
# SLIDE 12 — SYSTEM ARCHITECTURE (diagram)
# ===========================================================================
diagram_slide("Design", "System Architecture",
    "07_system_architecture.png",
    "Figure 3: Layered architecture — presentation (Flutter), logic (Providers/Services), data (Firebase).",
    "This is the system architecture, organized in three layers. The top presentation layer is the "
    "Flutter UI — all the screens. The middle layer is the application logic — my Providers and Services "
    "handle state and talk to the backend. The bottom data layer is Firebase: Auth, Firestore, Storage, "
    "Cloud Functions and Messaging. Separating these layers keeps the UI independent of the backend and "
    "makes the code easier to maintain and test.")

# ===========================================================================
# SLIDE 13 — USE CASE (diagram)
# ===========================================================================
diagram_slide("Requirements", "Use Case Diagram",
    "02_use_case_UML.png",
    "Figure 4: Two actors — User/Donor and Admin — and their interactions with the system.",
    "This use-case diagram captures what each actor can do. There are two actors. The User/Donor can "
    "register, log in, post and view blood requests, toggle donor availability, search donors, chat, and "
    "browse the emergency directory. The Admin manages the curated content — ambulances, organizations "
    "and volunteers. Both authenticate through Firebase Auth. This diagram was my functional-requirements "
    "blueprint for the whole build.")

# ===========================================================================
# SLIDE 14 — DFD LEVEL 0 (diagram)
# ===========================================================================
diagram_slide("Data Flow", "Data Flow Diagram — Level 0 (Context)",
    "03_DFD_level0_context.png",
    "Figure 5: The whole system as one process, with external entities and data flows.",
    "The context-level data flow diagram shows the entire Life Link system as a single process in the "
    "middle, with the external entities around it — the user, the admin, and Firebase services. The arrows "
    "show what data moves where: users send requests and profiles in, and receive notifications and "
    "matches out. It's the highest-level view before we break the system into detailed processes.")

# ===========================================================================
# SLIDE 15 — DFD LEVEL 1 (diagram)
# ===========================================================================
diagram_slide("Data Flow", "Data Flow Diagram — Level 1",
    "04_DFD_level1.png",
    "Figure 6: The system decomposed into its main processes and data stores.",
    "Level 1 zooms in and breaks the single process into the major sub-processes: authentication, "
    "request management, donor matching, notifications, chat and the directory. You can see how each "
    "process reads from and writes to the Firestore data stores — Users, Blood Requests, Chats and so on. "
    "This decomposition guided how I structured my services in code.")

# ===========================================================================
# SLIDE 16 — ERD (diagram)
# ===========================================================================
diagram_slide("Data Model", "Entity Relationship Diagram (ERD)",
    "01_ERD.png",
    "Figure 7: Core entities — User, Blood Request, Chat, Message, Notification, and the directory tables.",
    "The ERD shows the data model. The central entity is the User, which relates to Blood Requests (a "
    "user creates many requests), to Chats and Messages (users exchange messages), and to Notifications. "
    "Separate entities hold the directory data — Organizations, Ambulances and Volunteers. The "
    "relationships and keys here map directly to my Firestore collections.")

# ===========================================================================
# SLIDE 17 — DATABASE DESIGN (diagram)
# ===========================================================================
diagram_slide("Data Model", "Database Design — Firestore Collections",
    "08_database_design.png",
    "Figure 8: NoSQL document schema for each Firestore collection and sub-collection.",
    "Because Firestore is a NoSQL document database, this slide shows the actual collection schema rather "
    "than relational tables. The main collections are users, Blood_request, chats with a messages "
    "sub-collection, Ambulance, organizations, Volunteer, and a per-user notifications sub-collection. "
    "Each box lists the fields and their types. This is the concrete shape of the data my app reads and "
    "writes.")

# ===========================================================================
# SLIDE 18 — SECURITY
# ===========================================================================
s = add_slide(); rect(s, 0, 0, SW, SH, WHITE)
header(s, "Trust & Safety", "Security & Data Protection")
bullets(s, [
    "Authentication required: every read/write demands a signed-in user (request.auth != null).",
    "Data ownership: users can only edit or delete their own profile and their own requests.",
    "Chat privacy: only the two participants of a conversation can read or write its messages.",
    "Directory data (ambulances, organizations, volunteers) is read-only for normal users.",
    "Image uploads to Firebase Storage are size-limited and restricted to authenticated users.",
    "Rules are enforced server-side by Firestore & Storage — they cannot be bypassed by the client.",
], top=Inches(2.1), size=16, gap=13)
notes(s, "Because this is health-related data, security mattered. I wrote Firestore and Storage security "
         "rules that enforce four things. One: every operation requires authentication. Two: data "
         "ownership — you can only change your own profile and your own requests. Three: chat privacy — "
         "only the two people in a conversation can see its messages. Four: the curated directory is "
         "read-only for normal users. Crucially, these rules run on Google's servers, so even a tampered "
         "client cannot bypass them.")

# ===========================================================================
# SLIDE 19 — TESTING
# ===========================================================================
s = add_slide(); rect(s, 0, 0, SW, SH, WHITE)
header(s, "Quality", "Testing & Validation")
bullets(s, [
    "Functional testing: every feature manually tested against its requirement (register, post, match, chat).",
    "Matching tests: verified donors are notified only on exact blood-group + city match, creator excluded.",
    "Notification testing: confirmed push delivery in foreground and background on real devices.",
    "Cross-device testing: ran on multiple Android screen sizes (responsive UI via ScreenUtil).",
    "Security testing: attempted unauthorized reads/writes — correctly blocked by Firestore rules.",
    "Usability testing: informal feedback from peers to refine navigation and labels.",
], top=Inches(2.1), size=16, gap=12)
notes(s, "I validated the app on several levels. Functional testing checked every feature against its "
         "requirement. I specifically tested the matching logic to confirm that only correct-blood-group, "
         "same-city donors get notified and the creator is excluded. I verified notifications arrive in "
         "both foreground and background on real phones. I tested across different screen sizes thanks to "
         "the responsive layout, and I did security testing by trying unauthorized access, which the rules "
         "correctly blocked. Finally, peers used the app and their feedback helped me refine the UI.")

# ===========================================================================
# SLIDE 20 — RESULTS / IMPACT
# ===========================================================================
s = add_slide(); rect(s, 0, 0, SW, SH, WHITE)
header(s, "Outcome", "Results & Impact")
# stat chips
chip(s, Inches(0.85), Inches(2.15), Inches(2.7), "Core Features", "9+")
chip(s, Inches(3.75), Inches(2.15), Inches(2.7), "App Screens", "25+")
chip(s, Inches(6.65), Inches(2.15), Inches(2.7), "Platforms", "iOS + Android")
chip(s, Inches(9.55), Inches(2.15), Inches(2.7), "Match Time", "Seconds")
bullets(s, [
    "Delivered a fully working cross-platform app from a single Flutter codebase.",
    "Reduced donor discovery from many manual phone calls to one automatic, targeted broadcast.",
    "Real-time matching, chat and notifications all functioning on live Firebase infrastructure.",
    "Centralized emergency services (ambulances, blood banks, volunteers) into one accessible place.",
], top=Inches(3.85), size=16, gap=13)
notes(s, "So what did I actually achieve? A complete, working cross-platform app with over nine core "
         "features across more than twenty-five screens, running on both iOS and Android from one "
         "codebase. The headline result is the matching: what used to be dozens of manual phone calls is "
         "now a single automatic, targeted broadcast that reaches the right donors in seconds. Real-time "
         "chat and notifications are live on Firebase, and all the emergency services are now in one place.")

# ===========================================================================
# SLIDE 21 — CHALLENGES & LEARNINGS
# ===========================================================================
s = add_slide(); rect(s, 0, 0, SW, SH, WHITE)
header(s, "Reflection", "Challenges & Lessons Learned")
bullets(s, [
    "Cloud Functions & FCM: getting server-side matching and reliable push delivery right took iteration.",
    ("Learned to normalize data (case, whitespace) so matching is robust.", 1),
    "Real-time chat: managing Firestore streams without duplicate listeners or memory leaks.",
    ("Learned to cache subscriptions and dispose them properly.", 1),
    "State management: Provider taught me to keep UI and business logic cleanly separated.",
    "Security rules: writing rules that are strict yet don't break legitimate features.",
    "Overall: end-to-end ownership of a real product — design, build, test and document.",
], top=Inches(2.05), size=16, gap=10)
notes(s, "Every project teaches you something. The hardest part was the server side — getting Cloud "
         "Functions and push notifications to fire reliably, which taught me to normalize data so matching "
         "isn't broken by capitalization or spaces. Real-time chat taught me to manage Firestore "
         "subscriptions carefully to avoid leaks. Provider showed me the value of separating UI from "
         "logic, and writing security rules taught me to balance strictness with usability. Most "
         "importantly, I owned a real product end to end — from design through testing and documentation.")

# ===========================================================================
# SLIDE 22 — FUTURE WORK
# ===========================================================================
s = add_slide(); rect(s, 0, 0, SW, SH, WHITE)
header(s, "What's Next", "Future Enhancements")
bullets(s, [
    "GPS & live maps: show nearby donors and route to the hospital in real time.",
    "Donation history & gamification: badges and reminders to encourage repeat donation.",
    "Verified blood banks: integrate official inventory so users see real-time stock.",
    "Admin web dashboard: manage directory content and view analytics.",
    "Multi-language support to widen accessibility.",
    "AI-assisted urgency ranking to prioritize the most critical requests.",
    "Optional medical verification for added trust and safety.",
], top=Inches(2.05), size=16, gap=11)
notes(s, "There's a clear roadmap beyond the FYP. I'd add GPS and live maps to show nearby donors and "
         "routes. I'd add donation history with gamification to encourage repeat donors. Integrating "
         "verified blood banks would show real-time stock. An admin web dashboard would make managing "
         "content and viewing analytics easier. Beyond that: multi-language support, AI-based urgency "
         "ranking, and optional medical verification. These show the project has room to grow into a "
         "production-grade service.")

# ===========================================================================
# SLIDE 23 — CONCLUSION
# ===========================================================================
s = add_slide(); rect(s, 0, 0, SW, SH, WHITE)
header(s, "Wrap-up", "Conclusion")
bullets(s, [
    "Life Link tackles a real, life-critical problem: finding matching blood donors fast.",
    "It replaces slow manual calling with automatic, location-and-blood-group-aware matching.",
    "Built on a modern, scalable stack — Flutter for one cross-platform app, Firebase for a real-time serverless backend.",
    "All core objectives were met: requests, matching, notifications, chat, directory and security.",
    "The result is a practical platform with genuine potential to save lives.",
], top=Inches(2.2), size=18, gap=16)
notes(s, "To conclude: Life Link addresses a real and life-critical problem by replacing slow manual "
         "calling with automatic, location- and blood-group-aware matching. It's built on a modern, "
         "scalable stack — Flutter for a single cross-platform app and Firebase for a real-time serverless "
         "backend. I met all of my core objectives, and the result is a practical platform that genuinely "
         "has the potential to save lives. Thank you.")

# ===========================================================================
# SLIDE 24 — THANK YOU / Q&A
# ===========================================================================
s = add_slide(); rect(s, 0, 0, SW, SH, INK)
rect(s, 0, 0, Inches(0.45), SH, RED)
rect(s, 0, SH - Inches(0.18), SW, Inches(0.18), RED)
tb, tf = textbox(s, Inches(1.0), Inches(2.6), Inches(11.3), Inches(1.5))
set_run(tf.paragraphs[0].add_run(), "Thank You", 60, WHITE, bold=True)
tb, tf = textbox(s, Inches(1.05), Inches(3.9), Inches(11), Inches(0.7))
set_run(tf.paragraphs[0].add_run(), "Questions & Discussion", 26, ACCENT, bold=True)
if os.path.exists("assets_pres/uni_logo.png"):
    s.shapes.add_picture("assets_pres/uni_logo.png", Inches(10.45), Inches(0.6),
                         width=Inches(1.9), height=Inches(1.9))
tb, tf = textbox(s, Inches(1.05), Inches(5.0), Inches(11), Inches(1.4))
p = tf.paragraphs[0]
set_run(p.add_run(), "Life Link", 16, WHITE, bold=True)
set_run(p.add_run(), "  —  A Blood Donation & Emergency Response Platform", 16,
        RGBColor(0xC9,0xCC,0xD1))
p2 = tf.add_paragraph(); p2.space_before = Pt(8)
set_run(p2.add_run(), "Ubaid Ullah  •  M. Saad Awan  •  Umar Rashid", 14, WHITE)
p3 = tf.add_paragraph(); p3.space_before = Pt(4)
set_run(p3.add_run(),
        "Supervisor: Engr. Waseem Ullah Khan  •  ICS/IT, The University of Agriculture, Peshawar",
        13, SLATE)
notes(s, "That brings me to the end of my presentation. Thank you for your time and attention. "
         "I'd be happy to take any questions, and I can give a live demo of any feature you'd like to "
         "see. Tip: be ready for common questions — why Flutter over native, how matching scales with "
         "many users, how you handle privacy of medical data, and what happens if no donor matches.")

# ---- Save -----------------------------------------------------------------
out = "Life_Link_FYP_Presentation.pptx"
prs.save(out)
print(f"Saved {out} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
